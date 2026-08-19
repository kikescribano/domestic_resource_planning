"""Inventario de la plantilla de Slidesgo: qué hay en cada diapositiva y con qué
identificador se rellena.

Es el paso previo a escribir un generador. Sin esto se acaba adivinando
identificadores de forma, y una forma equivocada no da error: escribe el texto
en un sitio que no se ve.

    python inspect-template.py            # resumen de las 42 diapositivas
    python inspect-template.py 9 10 21    # detalle de las que se indiquen

En el detalle, cada línea de texto lleva el **número de caracteres del relleno
original**. Ese número es el presupuesto: la plantilla está compuesta en inglés y
el castellano es más largo, así que pasarse de largo no es una opción estética
sino un desbordamiento. Ver `qa-deck.py`, que lo comprueba.
"""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
from slidesgo_deck import TEMPLATE  # noqa: E402

from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402

sys.stdout.reconfigure(encoding="utf-8")


def inches(value) -> float:
    return round(Emu(value).inches, 2) if value is not None else 0.0


def walk(shapes, depth: int = 0):
    for shape in shapes:
        yield depth, shape
        if shape.shape_type == 6:  # GROUP
            yield from walk(shape.shapes, depth + 1)


def summary(prs: Presentation) -> None:
    print(f"{len(prs.slides)} diapositivas · lienzo "
          f"{inches(prs.slide_width)} x {inches(prs.slide_height)} pulgadas\n")
    for number, slide in enumerate(prs.slides, 1):
        titles = [s.text_frame.paragraphs[0].text
                  for s in slide.shapes
                  if s.has_text_frame and s.text_frame.text.strip()]
        print(f"{number:>2} {slide.slide_layout.name:<30} "
              f"{' | '.join(titles[:3])[:78]}")


def detail(prs: Presentation, number: int) -> None:
    slide = prs.slides[number - 1]
    print(f"\n===== DIAPOSITIVA {number} · patrón {slide.slide_layout.name} =====")
    for depth, shape in walk(slide.shapes):
        sangria = "  " * depth
        cabecera = (f"{sangria}id={shape.shape_id:<5} {str(shape.shape_type):<16} "
                    f"@({inches(shape.left)},{inches(shape.top)}) "
                    f"{inches(shape.width)}x{inches(shape.height)}")
        if shape.has_text_frame and shape.text_frame.text.strip():
            print(cabecera)
            for parrafo in shape.text_frame.paragraphs:
                if not parrafo.text.strip():
                    continue
                fuente = parrafo.runs[0].font if parrafo.runs else None
                medida = fuente.size.pt if fuente is not None and fuente.size else "heredado"
                print(f"{sangria}      nivel={parrafo.level} pt={medida} "
                      f"[{len(parrafo.text)} car] {parrafo.text!r}")
        elif shape.has_table:
            tabla = shape.table
            print(f"{cabecera}  TABLA {len(tabla.rows)}x{len(tabla.columns)}")
            for fila in tabla.rows:
                print(f"{sangria}      {[celda.text for celda in fila.cells]}")
        elif shape.shape_type == 13:  # PICTURE
            print(f"{cabecera}  IMAGEN {shape.image.filename or shape.image.content_type} "
                  f"({shape.image.size[0]}x{shape.image.size[1]}, "
                  f"{len(shape.image.blob) // 1024} kB)")


def main() -> None:
    prs = Presentation(str(TEMPLATE))
    numeros = [int(a) for a in sys.argv[1:]]
    if not numeros:
        summary(prs)
        return
    for numero in numeros:
        detail(prs, numero)


if __name__ == "__main__":
    main()
