"""Verifica una presentación derivada de la plantilla de Slidesgo.

Comprueba cuatro cosas, y las cuatro han fallado alguna vez:

1. **Atribución.** La licencia de cuenta gratuita obliga a conservar la
   diapositiva de agradecimiento. Sin ella la presentación no se entrega.
2. **Relleno olvidado.** Una forma que no se tocó sigue hablando de Venus, de
   Marte o de John Doe, y no da ningún error: se ve perfecta hasta que alguien
   la lee.
3. **Presupuesto de texto.** La plantilla está compuesta en inglés y el
   castellano ocupa más. Se compara cada texto con el relleno original de esa
   misma forma; pasarse mucho es el aviso barato de un desbordamiento.
4. **Encaje.** Se mide cuánto ocupa el texto al partirse en líneas, resolviendo
   el cuerpo que hereda del patrón, y se compara con **dos** referencias: la
   caja y lo que ocupaba el relleno original de esa misma forma.

Sobre la medida, que tiene una trampa: la plantilla incrusta Montserrat, Roboto
y Bebas Neue, pero las incrusta como **EOT comprimido con MicroType Express**, y
eso no lo abre ninguna librería de las que hay aquí. Así que se mide con una
tipografía sustituta del sistema. Para que la sustitución no invente
desbordamientos, **la comparación se calibra sola**: el mismo texto original se
mide con la misma fuente equivocada, y lo que se compara son las dos medidas
entre sí. Si el relleno ya «no cabía» según la sustituta, la referencia pasa a
ser el relleno y no la caja.

    pip install python-pptx Pillow
    python qa-deck.py deck.pptx

Sale con código 1 si hay algo que impide entregar: falta la atribución, queda
relleno de la plantilla o un texto se sale de su caja.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

sys.path.insert(0, str(Path(__file__).resolve().parent))
from slidesgo_deck import ATTRIBUTION_SLIDE, TEMPLATE  # noqa: E402

sys.stdout.reconfigure(encoding="utf-8")

A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P_NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

# Palabras del relleno de la plantilla. Si sobreviven, es que una forma se quedó
# sin tocar. «Slidesgo», «Freepik» y «Flaticon» no están: viven en los créditos
# del patrón de la diapositiva de agradecimiento y ahí tienen que seguir.
RELLENO = ("Mercury", "Venus", "Mars", "Jupiter", "Saturn", "Neptune",
           "John Doe", "Susan Bones", "Tommy Lean", "Lorem ipsum",
           "You can describe", "You can speak about", "Here is where",
           "your presentation begins", "AWESOME WORDS", "PRODUCT 1",
           "Big numbers catch", "planet")

EXCESO_AVISO = 1.25
EXCESO_ERROR = 1.60


# ---------------------------------------------------------------------------
# Tipografías: sustitutas del sistema, porque las incrustadas no se pueden leer
# ---------------------------------------------------------------------------

FONTDIR = Path(r"C:\Windows\Fonts")
SUSTITUTAS = {
    "montserrat": ("segoeui.ttf", "seguisb.ttf"),
    "roboto": ("arial.ttf", "arialbd.ttf"),
    "proxima": ("arial.ttf", "arialbd.ttf"),
    "bebas": ("arialbd.ttf", "arialbd.ttf"),
    "amatic": ("arial.ttf", "arialbd.ttf"),
}
_cache: dict[tuple[str, bool, int], object] = {}


def load_font(nombre: str, negrita: bool, puntos: float, dpi: int = 96):
    from PIL import ImageFont

    pixeles = max(1, int(round(puntos * dpi / 72.0)))
    clave = (nombre.lower(), negrita, pixeles)
    if clave in _cache:
        return _cache[clave]
    familia = next((v for k, v in SUSTITUTAS.items() if k in nombre.lower()),
                   ("arial.ttf", "arialbd.ttf"))
    ruta = FONTDIR / familia[1 if negrita else 0]
    if not ruta.exists():
        ruta = FONTDIR / "arial.ttf"
    if not ruta.exists():
        return None
    fuente = ImageFont.truetype(str(ruta), pixeles)
    _cache[clave] = fuente
    return fuente


# ---------------------------------------------------------------------------
# Formato heredado: cuerpo y tipografía efectivos de un párrafo
# ---------------------------------------------------------------------------

def _lst_styles(shape, slide):
    """Los `lstStyle` que gobiernan una forma, del más concreto al más general."""
    estilos = []
    cuerpo = shape._element.find(f"{P_NS}txBody")
    if cuerpo is not None:
        propio = cuerpo.find(f"{A_NS}lstStyle")
        if propio is not None:
            estilos.append(propio)
    if not shape.is_placeholder:
        return estilos

    indice = shape.placeholder_format.idx
    patron = slide.slide_layout
    for contenedor in (patron, patron.slide_master):
        for marcador in contenedor.placeholders:
            if marcador.placeholder_format.idx != indice:
                continue
            cuerpo = marcador._element.find(f"{P_NS}txBody")
            if cuerpo is None:
                continue
            estilo = cuerpo.find(f"{A_NS}lstStyle")
            if estilo is not None:
                estilos.append(estilo)
    maestro = patron.slide_master._element.find(f"{P_NS}txStyles")
    if maestro is not None:
        tipo = str(shape.placeholder_format.type)
        etiqueta = ("titleStyle" if "TITLE" in tipo else
                    "bodyStyle" if "BODY" in tipo or "SUBTITLE" in tipo else "otherStyle")
        estilo = maestro.find(f"{P_NS}{etiqueta}")
        if estilo is not None:
            estilos.append(estilo)
    return estilos


def effective(shape, slide, paragraph, run) -> tuple[float, str, bool]:
    """(cuerpo en puntos, tipografía, negrita) efectivos de un run."""
    puntos = run.font.size.pt if run.font.size else None
    nombre = run.font.name
    negrita = bool(run.font.bold)

    nivel = f"{A_NS}lvl{paragraph.level + 1}pPr"
    for estilo in _lst_styles(shape, slide):
        propiedades = estilo.find(nivel)
        if propiedades is None:
            continue
        defecto = propiedades.find(f"{A_NS}defRPr")
        if defecto is None:
            continue
        if puntos is None and defecto.get("sz"):
            puntos = int(defecto.get("sz")) / 100.0
        if nombre is None:
            latina = defecto.find(f"{A_NS}latin")
            if latina is not None and latina.get("typeface"):
                nombre = latina.get("typeface")
        if puntos is not None and nombre is not None:
            break
    return puntos or 18.0, nombre or "Roboto", negrita


def line_spacing(paragraph, puntos: float) -> float:
    espaciado = paragraph.line_spacing
    if espaciado is None:
        return puntos * 1.2
    if isinstance(espaciado, (int, float)):
        return puntos * float(espaciado)
    return espaciado.pt  # ya venía en puntos


# ---------------------------------------------------------------------------
# Comprobaciones
# ---------------------------------------------------------------------------

def template_lengths() -> dict[int, tuple[str, str]]:
    """{id de forma: (texto original, título de su diapositiva)} de la plantilla."""
    prs = Presentation(str(TEMPLATE))
    inventario: dict[int, tuple[str, str]] = {}

    def recorrer(shapes, etiqueta):
        for shape in shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                inventario.setdefault(shape.shape_id, (shape.text_frame.text, etiqueta))
            if shape.shape_type == 6:
                recorrer(shape.shapes, etiqueta)

    for numero, slide in enumerate(prs.slides, 1):
        recorrer(slide.shapes, f"plantilla {numero}")
    return inventario


def walk(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:
            yield from shape.shapes


def check(path: Path) -> int:
    prs = Presentation(str(path))
    original = template_lengths()

    errores: list[str] = []
    avisos: list[str] = []
    atribucion = False

    print(f"{path.name}: {len(prs.slides)} diapositivas\n")

    for numero, slide in enumerate(prs.slides, 1):
        titulo = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                titulo = shape.text_frame.paragraphs[0].text
                break
        print(f"{numero:>2} {slide.slide_layout.name:<24} {titulo[:60]}")

        for shape in walk(slide.shapes):
            if shape.shape_id in (729, 731) or "THANKS" in (
                    shape.text_frame.text if shape.has_text_frame else ""):
                atribucion = True
            if not shape.has_text_frame:
                continue
            texto = shape.text_frame.text
            if not texto.strip():
                continue

            for palabra in RELLENO:
                if palabra.lower() in texto.lower():
                    errores.append(
                        f"diapositiva {numero}, id={shape.shape_id}: queda relleno de "
                        f"la plantilla ({palabra!r}) — {texto[:60]!r}")
                    break

            referencia = original.get(shape.shape_id)
            if referencia and len(referencia[0]) > 24:
                razon = len(texto) / len(referencia[0])
                if razon >= EXCESO_ERROR:
                    errores.append(
                        f"diapositiva {numero}, id={shape.shape_id}: {len(texto)} car "
                        f"frente a {len(referencia[0])} del relleno (x{razon:.1f})")
                elif razon >= EXCESO_AVISO:
                    avisos.append(
                        f"diapositiva {numero}, id={shape.shape_id}: {len(texto)} car "
                        f"frente a {len(referencia[0])} del relleno (x{razon:.1f})")

            desborde = overflow(shape, slide, referencia[0] if referencia else None)
            if desborde:
                errores.append(f"diapositiva {numero}, id={shape.shape_id}: {desborde}")

    if not atribucion:
        errores.append(
            f"falta la diapositiva de agradecimiento (la {ATTRIBUTION_SLIDE} de la "
            f"plantilla): la licencia de cuenta gratuita obliga a conservarla")

    print()
    for aviso in avisos:
        print(f"AVISO  {aviso}")
    for error in errores:
        print(f"ERROR  {error}")
    if not avisos and not errores:
        print("sin incidencias")
    return 1 if errores else 0


def _alto(lineas: list[str], formatos: list[tuple], ancho: float, lapiz) -> float:
    """Alto en píxeles que ocupan unas líneas de texto con unos formatos dados."""
    total = 0.0
    for indice, linea in enumerate(lineas):
        puntos, nombre, negrita, espaciado = formatos[min(indice, len(formatos) - 1)]
        if not linea.strip():
            total += espaciado * 96 / 72.0
            continue
        fuente = load_font(nombre, negrita, puntos)
        if fuente is None:
            return 0.0
        cuenta, actual = 1, ""
        for palabra in linea.split():
            prueba = f"{actual} {palabra}".strip()
            if lapiz.textlength(prueba, font=fuente) <= ancho or not actual:
                actual = prueba
            else:
                cuenta += 1
                actual = palabra
        total += cuenta * espaciado * 96 / 72.0
    return total


def overflow(shape, slide, original: str | None) -> str | None:
    """Estima si el texto se sale de su caja, calibrando con el relleno original."""
    from PIL import Image, ImageDraw

    marco = shape.text_frame
    if shape.width is None or shape.height is None or marco.word_wrap is False:
        return None
    ancho = (Emu(shape.width).inches
             - Emu(marco.margin_left).inches - Emu(marco.margin_right).inches) * 96
    alto = (Emu(shape.height).inches
            - Emu(marco.margin_top).inches - Emu(marco.margin_bottom).inches) * 96
    if ancho <= 4 or alto <= 4:
        return None

    formatos = []
    for parrafo in marco.paragraphs:
        run = parrafo.runs[0] if parrafo.runs else None
        if run is None:
            continue
        puntos, nombre, negrita = effective(shape, slide, parrafo, run)
        formatos.append((puntos, nombre, negrita, line_spacing(parrafo, puntos)))
    if not formatos:
        return None

    lapiz = ImageDraw.Draw(Image.new("RGB", (8, 8)))
    actual = _alto([p.text for p in marco.paragraphs], formatos, ancho, lapiz)
    referencia = alto
    if original:
        del_relleno = _alto(original.split("\n"), formatos, ancho, lapiz)
        referencia = max(alto, del_relleno)  # si el relleno ya «no cabía», manda él

    if actual > referencia * 1.06:
        return (f"el texto ocupa ~{actual / 96:.2f}\" donde caben "
                f"{referencia / 96:.2f}\" — {marco.text[:50]!r}")
    return None


if __name__ == "__main__":
    if len(sys.argv) != 2:
        raise SystemExit(__doc__)
    raise SystemExit(check(Path(sys.argv[1])))
