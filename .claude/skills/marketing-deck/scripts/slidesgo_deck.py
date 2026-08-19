"""Construye presentaciones derivadas de «Pitch Deck Minitheme» de Slidesgo.

La idea de fondo es la que conserva el look and feel: **no se dibuja nada, se
parte de la plantilla**. Una presentación de DRP es una *selección* de las
diapositivas de la plantilla, reordenadas y con el texto sustituido. Todo lo que
hace que se reconozca —el negro, el verde `#9CFC34`, los círculos a mano alzada,
los bordes de papel rasgado, las ilustraciones y **las tipografías incrustadas**—
viaja con la diapositiva sin que haya que reproducirlo.

    from slidesgo_deck import Deck

    deck = Deck(titulo="DRP · documento comercial")
    portada = deck.use(1)
    portada.text(165, "DRP")
    ...
    deck.save("salida.pptx")

Tres cosas que este módulo impone y no son negociables:

- **La diapositiva de agradecimiento va siempre.** La plantilla se descargó con
  cuenta gratuita y esa licencia exige conservarla; `save()` falla si no está.
- **El lienzo es el de la plantilla** —10 x 5,625 pulgadas—, porque las
  diapositivas se reutilizan tal cual.
- **El texto sustituye, no reformatea.** Se conserva el formato del párrafo y de
  su primer run, que es de donde salen tipografía, cuerpo y color heredados.

Requiere `pip install python-pptx`.
"""

from __future__ import annotations

import copy
import shutil
from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# La plantilla de referencia, versionada en el área de marketing.
TEMPLATE = (Path(__file__).resolve().parents[4]
            / "docs" / "common" / "marketing" / "references"
            / "pitch-deck-minitheme-slidesgo.pptx")

# La diapositiva «THANKS». La licencia de usuario gratuito obliga a conservarla
# en cualquier presentación derivada: es la atribución a Slidesgo. Los créditos
# —Slidesgo, Flaticon y Freepik— no están en la diapositiva sino en su patrón
# `BLANK_1`, así que llegan solos con ella.
ATTRIBUTION_SLIDE = 21

_R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


class DeckError(RuntimeError):
    """Error de construcción: identificador que no existe, atribución ausente."""


def _set_paragraph(paragraph, text: str, idioma: str = "es-ES") -> None:
    """Escribe `text` en un párrafo conservando el formato de su primer run.

    Con una excepción: si ese run era un **enlace**, se le quita el formato
    entero. La plantilla marca sus enlaces en verde, negrita y subrayado, y ese
    formato vive en el run, no en el enlace: reutilizarlo tal cual deja una línea
    resaltada en mitad de una lista, sin ninguna razón visible.
    """
    runs = paragraph.runs
    if not runs:
        # Sin runs no hay formato propio que conservar: el que aparezca lo hereda
        # del párrafo, del marcador de posición del patrón y del tema.
        paragraph.add_run().text = text
        return

    propiedades = runs[0]._r.find(f"{_A_NS}rPr")
    if propiedades is not None:
        if propiedades.find(f"{_A_NS}hlinkClick") is not None:
            for hijo in list(propiedades):
                propiedades.remove(hijo)
            for atributo in list(propiedades.attrib):
                if atributo != "lang":
                    propiedades.attrib.pop(atributo)
        propiedades.set("lang", idioma)
    runs[0].text = text
    for sobrante in runs[1:]:
        sobrante._r.getparent().remove(sobrante._r)


def _reset_autofit(text_frame) -> None:
    """Quita la reducción de cuerpo que la plantilla calculó para *su* texto.

    Google Slides deja escrito en `normAutofit` cuánto encogió la letra para que
    cupiera el relleno original. Ese factor no se recalcula al sustituir el
    texto: se queda, y el resultado es una diapositiva con seis palabras escritas
    en el cuerpo que hacía falta para trece líneas —o dos tarjetas iguales con
    dos tamaños de letra distintos, que es como se descubrió—.
    """
    cuerpo = text_frame._txBody.find(f"{_A_NS}bodyPr")
    if cuerpo is None:
        return
    ajuste = cuerpo.find(f"{_A_NS}normAutofit")
    if ajuste is None:
        return
    ajuste.attrib.pop("fontScale", None)
    ajuste.attrib.pop("lnSpcReduction", None)


def _remap_relationship_ids(element, mapping: dict[str, str]) -> None:
    """Reescribe los `r:id`/`r:embed` de un árbol copiado a otra diapositiva."""
    for node in element.iter():
        for nombre, valor in list(node.attrib.items()):
            if nombre.startswith(_R_NS) and valor in mapping:
                node.set(nombre, mapping[valor])


class Slide:
    """Una diapositiva ya elegida, lista para que se le sustituya el contenido."""

    def __init__(self, slide, origen: int):
        self._slide = slide
        self.origen = origen

    # -- localización de formas ------------------------------------------------

    def _shape(self, shape_id: int):
        for shape in self._walk(self._slide.shapes):
            if shape.shape_id == shape_id:
                return shape
        raise DeckError(
            f"la diapositiva {self.origen} no tiene ninguna forma con id={shape_id}; "
            f"míralo con `inspect-template.py {self.origen}`")

    @classmethod
    def _walk(cls, shapes) -> Iterable:
        for shape in shapes:
            yield shape
            if shape.shape_type == 6:  # GROUP
                yield from cls._walk(shape.shapes)

    # -- contenido -------------------------------------------------------------

    def text(self, shape_id: int, *lines: str) -> "Slide":
        """Sustituye el texto de una forma, una línea por párrafo.

        Con menos líneas que párrafos, sobran párrafos y se borran; con más, el
        último se clona tantas veces como haga falta, de modo que una lista de
        viñetas crece conservando su nivel y su viñeta.
        """
        if not lines:
            raise DeckError(f"id={shape_id}: no se pasó ninguna línea")
        marco = self._shape(shape_id).text_frame
        if not marco.paragraphs:
            raise DeckError(f"id={shape_id}: la forma no tiene párrafos que reutilizar")

        while len(marco.paragraphs) < len(lines):
            ultimo = marco.paragraphs[-1]._p
            ultimo.addnext(copy.deepcopy(ultimo))
        for parrafo, linea in zip(marco.paragraphs, lines):
            _set_paragraph(parrafo, linea)
        for sobrante in list(marco.paragraphs)[len(lines):]:
            sobrante._p.getparent().remove(sobrante._p)
        _reset_autofit(marco)
        return self

    def table(self, shape_id: int, rows: Sequence[Sequence[str]]) -> "Slide":
        """Rellena una tabla celda a celda, conservando el formato de cada una."""
        tabla = self._shape(shape_id).table
        if len(rows) > len(tabla.rows) or any(len(f) > len(tabla.columns) for f in rows):
            raise DeckError(
                f"id={shape_id}: la tabla es de {len(tabla.rows)}x{len(tabla.columns)} "
                f"y no cabe el contenido que se le pasa")
        for indice_fila, fila in enumerate(rows):
            for indice_columna, valor in enumerate(fila):
                marco = tabla.cell(indice_fila, indice_columna).text_frame
                _set_paragraph(marco.paragraphs[0], valor)
                for sobrante in list(marco.paragraphs)[1:]:
                    sobrante._p.getparent().remove(sobrante._p)
        return self

    def drop(self, *shape_ids: int) -> "Slide":
        """Elimina formas: la tarjeta que sobra cuando hay cinco cosas y seis huecos."""
        for shape_id in shape_ids:
            elemento = self._shape(shape_id)._element
            elemento.getparent().remove(elemento)
        return self

    def notes(self, text: str) -> "Slide":
        """Anota de dónde sale el contenido. Es lo que permite auditarlo después."""
        self._slide.notes_slide.notes_text_frame.text = text
        return self

    def _purge_dead_links(self) -> None:
        """Suelta los enlaces que quedaron sin dueño al sustituir el texto.

        El índice de la plantilla enlaza a sus propias diapositivas de recursos.
        Al reescribirlo, el enlace desaparece del texto pero **la relación sigue
        ahí**, y con ella las diapositivas enlazadas: se quedan dentro del
        fichero sin salir en la presentación —tres megabytes de GIF invisibles—
        porque nada las hace inalcanzables.
        """
        parte = self._slide.part
        for rId in list(parte.rels):
            if parte.rels[rId].reltype in (RT.HYPERLINK, RT.SLIDE):
                parte.drop_rel(rId)  # solo cae si ya no lo referencia nadie


class Deck:
    """Una presentación derivada: se eligen diapositivas y se les pone contenido."""

    def __init__(self, titulo: str, autor: str = "DRP · Domestic Resource Planning",
                 template: Path | str = TEMPLATE):
        self._template = Path(template)
        if not self._template.exists():
            raise DeckError(f"no encuentro la plantilla en {self._template}")
        self._prs = Presentation(str(self._template))
        self._origen = list(self._prs.slides)
        self._elegidas: list[Slide] = []
        self._usadas: set[int] = set()
        self._titulo = titulo
        self._autor = autor

    # -- selección -------------------------------------------------------------

    def use(self, numero: int) -> Slide:
        """Añade la diapositiva `numero` de la plantilla al final de la presentación.

        Repetir una composición está permitido: la segunda vez y las siguientes se
        trabaja sobre un duplicado, no sobre la misma diapositiva.
        """
        if not 1 <= numero <= len(self._origen):
            raise DeckError(f"la plantilla no tiene diapositiva {numero}")
        original = self._origen[numero - 1]
        slide = original if numero not in self._usadas else self._duplicate(original)
        self._usadas.add(numero)
        elegida = Slide(slide, numero)
        self._elegidas.append(elegida)
        return elegida

    def _duplicate(self, source):
        destino = self._prs.slides.add_slide(source.slide_layout)
        for shape in list(destino.shapes):
            shape._element.getparent().remove(shape._element)

        equivalencias: dict[str, str] = {}
        for rId, rel in source.part.rels.items():
            if rel.reltype == RT.SLIDE_LAYOUT:
                continue
            equivalencias[rId] = (
                destino.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
                if rel.is_external
                else destino.part.rels.get_or_add(rel.reltype, rel.target_part))

        arbol = destino.shapes._spTree
        for shape in source.shapes:
            copia = copy.deepcopy(shape._element)
            _remap_relationship_ids(copia, equivalencias)
            arbol.insert_element_before(copia, "p:extLst")
        return destino

    # -- salida ----------------------------------------------------------------

    def save(self, destino: Path | str) -> Path:
        """Borra lo no elegido, ordena, comprueba la atribución y escribe el fichero."""
        if ATTRIBUTION_SLIDE not in self._usadas:
            raise DeckError(
                f"falta la diapositiva {ATTRIBUTION_SLIDE} («THANKS»): la plantilla se "
                f"descargó con cuenta gratuita y su licencia obliga a conservarla en "
                f"toda presentación derivada")

        for elegida in self._elegidas:
            elegida._purge_dead_links()

        elegidas = {id(slide._slide.part) for slide in self._elegidas}
        lista = self._prs.slides._sldIdLst
        for referencia, slide in zip(list(lista), list(self._prs.slides)):
            if id(slide.part) not in elegidas:
                lista.remove(referencia)
                self._prs.part.drop_rel(referencia.rId)

        por_parte = {id(slide.part): referencia
                     for referencia, slide in zip(list(lista), list(self._prs.slides))}
        for slide in self._elegidas:
            lista.append(por_parte[id(slide._slide.part)])

        propiedades = self._prs.core_properties
        propiedades.title = self._titulo
        propiedades.author = self._autor
        propiedades.comments = (
            "Derivado de la plantilla «Pitch Deck Minitheme» de Slidesgo "
            "(cuenta gratuita: se conserva la diapositiva de agradecimiento).")

        destino = Path(destino)
        destino.parent.mkdir(parents=True, exist_ok=True)
        self._prs.save(str(destino))
        return destino


def copy_template(destino: Path | str) -> Path:
    """Copia la plantilla en blanco. Útil para trastear sin tocar la referencia."""
    destino = Path(destino)
    shutil.copyfile(TEMPLATE, destino)
    return destino
