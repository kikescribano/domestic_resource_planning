"""Previsualizador de .pptx para el QA visual, cuando no hay PowerPoint ni LibreOffice.

Dibuja las formas y el texto de cada diapositiva con las fuentes reales del
sistema, así que las medidas de encaje son fiables aunque el resultado no sea un
render fiel de PowerPoint. Reporta el texto que no cabe en su caja y escribe un
PNG por diapositiva para revisarlas a ojo.

    pip install python-pptx Pillow
    python preview-pptx.py deck.pptx qa/

Asume las fuentes de Windows (FONTDIR); en otro sistema hay que ajustar el mapa.
Ver SKILL-001-readme-to-deck.md para el procedimiento completo.
"""
import sys, os
from pptx import Presentation
from pptx.util import Emu, Length
from PIL import Image, ImageDraw, ImageFont

DPI = 110
FONTDIR = r"C:\Windows\Fonts"
FONTMAP = {
    ("Calibri", False): "calibri.ttf", ("Calibri", True): "calibrib.ttf",
    ("Cambria", False): "cambria.ttc", ("Cambria", True): "cambriab.ttf",
    ("Courier New", False): "cour.ttf", ("Courier New", True): "courbd.ttf",
    ("Arial", False): "arial.ttf", ("Arial", True): "arialbd.ttf",
}
_cache = {}


def font(name, bold, size_pt, italic=False):
    key = (name, bold, round(size_pt, 1), italic)
    if key in _cache:
        return _cache[key]
    fn = FONTMAP.get((name, bold)) or FONTMAP.get((name, False)) or "calibri.ttf"
    if italic:
        alt = {"calibri.ttf": "calibrii.ttf", "calibrib.ttf": "calibriz.ttf",
               "cambria.ttc": "cambriai.ttf", "cambriab.ttf": "cambriaz.ttf"}.get(fn)
        if alt and os.path.exists(os.path.join(FONTDIR, alt)):
            fn = alt
    px = max(1, int(round(size_pt * DPI / 72.0)))
    f = ImageFont.truetype(os.path.join(FONTDIR, fn), px)
    _cache[key] = f
    return f


def emu_px(v):
    return int(round(Emu(v).inches * DPI))


def rgb_of(color, default=None):
    try:
        if color and color.type is not None and color.rgb is not None:
            return "#" + str(color.rgb)
    except Exception:
        pass
    return default


def wrap(text, fnt, maxw, draw):
    out = []
    for hard in text.split("\n"):
        words, line = hard.split(" "), ""
        for w in words:
            t = (line + " " + w).strip()
            if draw.textlength(t, font=fnt) <= maxw or not line:
                line = t
            else:
                out.append(line)
                line = w
        out.append(line)
    return out


def render(path, outdir):
    prs = Presentation(path)
    SW, SH = emu_px(prs.slide_width), emu_px(prs.slide_height)
    problems = []
    files = []
    for idx, slide in enumerate(prs.slides, 1):
        img = Image.new("RGB", (SW, SH), "white")
        d = ImageDraw.Draw(img, "RGBA")
        bg = None
        try:
            bg = rgb_of(slide.background.fill.fore_color)
        except Exception:
            pass
        if bg:
            d.rectangle([0, 0, SW, SH], fill=bg)
        for shp in slide.shapes:
            x, y = emu_px(shp.left), emu_px(shp.top)
            w, h = emu_px(shp.width), emu_px(shp.height)
            st = str(getattr(shp, "shape_type", ""))
            ast = ""
            try:
                ast = str(shp.auto_shape_type)
            except Exception:
                pass
            fill = None
            try:
                fill = rgb_of(shp.fill.fore_color)
            except Exception:
                pass
            line = None
            try:
                line = rgb_of(shp.line.color)
            except Exception:
                pass
            if "GRAPHIC_FRAME" in st or "CHART" in st:
                d.rectangle([x, y, x + w, y + h], outline="#999999", width=2)
                d.text((x + 10, y + 10), "[CHART]", font=font("Calibri", True, 12), fill="#666666")
                continue
            if "OVAL" in ast:
                d.ellipse([x, y, x + w, y + h], fill=fill, outline=line)
            elif "LINE" in ast and h <= 2 or ("LINE" in ast and w <= 2):
                d.line([x, y, x + w, y + h], fill=line or "#999999", width=2)
            elif "ARROW" in ast:
                d.rectangle([x, y, x + w, y + h], fill=fill, outline=line)
            elif fill or line:
                d.rounded_rectangle([x, y, x + w, y + h], radius=8, fill=fill, outline=line)
            if not shp.has_text_frame:
                continue
            tf = shp.text_frame
            ml = emu_px(tf.margin_left); mr = emu_px(tf.margin_right)
            mt = emu_px(tf.margin_top); mb = emu_px(tf.margin_bottom)
            bx, by = x + ml, y + mt
            bw, bh = max(2, w - ml - mr), max(2, h - mt - mb)
            blocks = []
            for p in tf.paragraphs:
                runs = [r for r in p.runs]
                if not runs:
                    blocks.append((None, "", 0, 0, None, p))
                    continue
                txt = "".join(r.text for r in runs)
                r0 = runs[0]
                size = (r0.font.size.pt if r0.font.size else 18)
                name = r0.font.name or "Calibri"
                bold = bool(r0.font.bold)
                ital = bool(r0.font.italic)
                col = rgb_of(r0.font.color, "#000000")
                f = font(name, bold, size, ital)
                ls = p.line_spacing
                if ls is None:
                    lh_pt = size * 1.22
                elif isinstance(ls, Length):
                    lh_pt = ls.pt
                elif isinstance(ls, (int, float)) and ls < 5:
                    lh_pt = size * float(ls)
                else:
                    lh_pt = float(ls)
                lh_px = int(round(lh_pt * DPI / 72.0))
                sa = p.space_after.pt if p.space_after else 0
                bullet = txt.startswith("\u2022")
                for ln in wrap(txt, f, bw - (14 if bullet else 0), d):
                    blocks.append((f, ln, lh_px, size, col, p))
                blocks.append((None, "", int(sa * DPI / 72.0), 0, None, p))
            total = sum(b[2] for b in blocks)
            anchor = str(tf.vertical_anchor)
            if "MIDDLE" in anchor:
                cy = by + max(0, (bh - total) // 2)
            elif "BOTTOM" in anchor:
                cy = by + max(0, bh - total)
            else:
                cy = by
            if total > bh + 3:
                problems.append(f"slide {idx}: TEXT OVERFLOW {total-bh}px  box=({shp.left/914400:.2f},{shp.top/914400:.2f},{shp.width/914400:.2f}x{shp.height/914400:.2f})  '{tf.text[:60]}'")
            for f, ln, lh_px, size, col, p in blocks:
                if f is None:
                    cy += lh_px
                    continue
                al = str(p.alignment)
                tw = d.textlength(ln, font=f)
                if "CENTER" in al:
                    tx = bx + (bw - tw) / 2
                elif "RIGHT" in al:
                    tx = bx + bw - tw
                else:
                    tx = bx
                d.text((tx, cy + (lh_px - size * DPI / 72.0) / 2), ln, font=f, fill=col or "#000000")
                cy += lh_px
        out = os.path.join(outdir, f"slide-{idx:02d}.png")
        img.save(out, quality=90)
        files.append(out)
    print("\n".join(problems) if problems else "no text-overflow issues detected")
    print("---")
    print("\n".join(files))


if __name__ == "__main__":
    os.makedirs(sys.argv[2], exist_ok=True)
    render(sys.argv[1], sys.argv[2])
